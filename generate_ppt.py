from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "VendSell.ai"
subtitle.text = "AI-Powered Retail Intel & Supply Chain Engine\nPitch Deck for AI Hackathon"

# Function to add slide
def add_slide(title_text, bullet_points):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = bullet_points[0]
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point

# Slide 2: The Problem
add_slide("The Problem in Modern Retail", [
    "Inefficient inventory management leads to severe stockouts and dead stock accumulation.",
    "Manual product cataloging and data entry is time-consuming and error-prone.",
    "Lack of real-time, data-driven insights across the supply chain.",
    "Communication gaps between Retailers (Selling Places) and Suppliers (Vendors)."
])

# Slide 3: Our Solution
add_slide("Our Solution: VendSell.ai", [
    "A state-of-the-art Supply Chain Platform connecting Retailers and Vendors seamlessly.",
    "Automated Inventory Optimization driven by machine learning.",
    "Predictive Analytics highlighting Expiry risks and Stockout dates.",
    "Multimodal OCR Vision for instant product catalog scanning and entry."
])

# Slide 4: AI & Predictive Analytics
add_slide("Core AI Capabilities", [
    "ABC/XYZ Classification: Prioritizes capital allocation based on revenue contribution and demand volatility.",
    "Stockout Risk Matrix: Predicts exact shortage dates using daily sales velocity and lead times.",
    "Ghost SKU Detection: Identifies dead stock (0 sales) for automated markdown strategies.",
    "EOQ & Safety Stock: AI calculates the exact Economic Order Quantity dynamically."
])

# Slide 5: The Conversational AI Assistant
add_slide("Retail LLM Assistant", [
    "Powered by Hugging Face Llama-3.1-8B-Instruct.",
    "Context-Aware: Directly queried against the Neon Postgres inventory database.",
    "Actionable: Ask 'Which products are running low?' or 'Create a liquidation plan'.",
    "Generates strategic, data-backed plans formatted perfectly in markdown."
])

# Slide 6: Architecture & Tech Stack
add_slide("Architecture & Technology Stack", [
    "Frontend: React 18, Vite, Material-UI, Recharts (Dark Glassmorphism UI).",
    "Backend Core: Node.js & Express transpiled to Vercel Serverless.",
    "AI Microservice: Python 3.10+, FastAPI (Serverless integration).",
    "Database: Neon Serverless Postgres for instant cold starts & pooling.",
    "Deployment: Fully optimized Vercel edge deployment."
])

# Slide 7: Business Impact
add_slide("Impact & Future Roadmap", [
    "Immediate ROI: Reduced holding costs and prevention of lost sales due to stockouts.",
    "Scalability: Serverless architecture ensures zero downtime and infinite scaling.",
    "Future Integration: Autonomous vendor reordering and dynamic AI pricing.",
    "Thank You! Let's revolutionize retail."
])

prs.save("VendSell.ai_Pitch_Deck.pptx")
print("Presentation generated successfully as VendSell.ai_Pitch_Deck.pptx")
